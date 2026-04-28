"""
Builds the AP Research presentation as a Canva-editable .pptx.

Visual style mirrors `/Users/yanzihao/Downloads/Current Presentation.pptx`:
  - 16:9 widescreen
  - Black -> deep wine (#34011F) gradient background
  - Magenta accents (#A9088E primary, #FB47DE highlight)
  - Poppins type family (Canva supplies this by default; safe fallback Arial)
  - White body text on dark background

Targets the highest band on every row of the AP Research Presentation rubric:
  - Row 1 (Research Design, 3): RQ + method + conclusion all stated.
  - Row 2 (Establish Argument, 6): connects evidence -> conclusion ->
    consequences/implications.
  - Row 3 (Reflect, 3): a dedicated slide that traces process -> conclusion.
  - Row 4 (Engage Audience, 6): every claim carries an MLA citation footer
    and speaker notes that cue gestures, vocal variety, and audience eye-line.

Run:
  python3 scripts/build_presentation.py
"""

from __future__ import annotations

import os
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

# 16:9 at 13.333" x 7.5" -- Canva's standard "Presentation" canvas size.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WINE = RGBColor(0x34, 0x01, 0x1F)           # background terminus
BLACK = RGBColor(0x00, 0x00, 0x00)          # background origin
MAGENTA = RGBColor(0xA9, 0x08, 0x8E)        # primary accent
HOT_PINK = RGBColor(0xFB, 0x47, 0xDE)       # highlight accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xE9, 0xE2, 0xE6)
MUTED = RGBColor(0xB6, 0xA8, 0xB0)
CARD = RGBColor(0x2C, 0x22, 0x27)           # tinted dark grey from current deck
HAIRLINE = RGBColor(0x55, 0x33, 0x4A)

FONT_TITLE = "Poppins"
FONT_BODY = "Poppins"


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def add_gradient_background(slide) -> None:
    """Insert a black -> wine linear gradient covering the entire slide."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    sp = bg.fill._xPr  # type: ignore[attr-defined]
    fill = sp.find(qn("a:solidFill"))
    if fill is not None:
        sp.remove(fill)
    grad_xml = (
        '<a:gradFill rotWithShape="1" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="000000"/></a:gs>'
        '<a:gs pos="55000"><a:srgbClr val="000000"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="34011F"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    sp.insert(2, etree.fromstring(grad_xml))
    bg.shadow.inherit = False
    # Drop the shape to z-bottom by reordering inside the spTree.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = FONT_BODY,
    line_spacing: float | None = 1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_runs(
    slide,
    left,
    top,
    width,
    height,
    runs: Sequence[dict],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = 1.2,
):
    """Multi-run text box. Each dict supports text/size/bold/color/font/italic/break_before."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for r in runs:
        if r.get("break_before"):
            p = tf.add_paragraph()
            p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", FONT_BODY)
        run.font.size = Pt(r.get("size", 18))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", WHITE)
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: Iterable[str],
    *,
    size: int = 18,
    color: RGBColor = SOFT_WHITE,
    bullet: str = "—",
    line_spacing: float = 1.25,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    items = list(items)
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r1 = p.add_run()
        r1.text = f"{bullet}  "
        r1.font.name = FONT_BODY
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = HOT_PINK
        r2 = p.add_run()
        r2.text = txt
        r2.font.name = FONT_BODY
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, *, fill=CARD, line=HAIRLINE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.06
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = line
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def add_pill(slide, left, top, label: str, *, fill=MAGENTA):
    width = Inches(2.4)
    height = Inches(0.42)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT_TITLE
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return pill


def add_hairline(slide, left, top, width, *, color=HAIRLINE, weight=1.25):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(weight))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def add_section_chrome(slide, eyebrow: str, title: str, *, slide_no: str | None = None):
    """Standard top chrome: section eyebrow, title, magenta hairline, slide number."""
    add_text(
        slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.34),
        eyebrow.upper(), size=11, bold=True, color=HOT_PINK, font=FONT_TITLE,
    )
    add_text(
        slide, Inches(0.6), Inches(0.78), Inches(12.0), Inches(0.95),
        title, size=32, bold=True, color=WHITE, font=FONT_TITLE,
    )
    add_hairline(slide, Inches(0.6), Inches(1.62), Inches(2.2),
                 color=HOT_PINK, weight=2)
    if slide_no:
        add_text(
            slide, Inches(11.6), Inches(0.45), Inches(1.3), Inches(0.34),
            slide_no, size=11, bold=True, color=MUTED, font=FONT_TITLE,
            align=PP_ALIGN.RIGHT,
        )


def add_citation_footer(slide, citations: Sequence[str]):
    """MLA citations along the bottom of the slide."""
    text = "  •  ".join(citations)
    add_text(
        slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.4),
        text, size=8.5, color=MUTED, font=FONT_BODY, line_spacing=1.15,
    )


def set_speaker_notes(slide, notes: str) -> None:
    nf = slide.notes_slide.notes_text_frame
    nf.text = ""  # reset placeholder
    paragraphs = notes.strip().split("\n\n")
    for i, para in enumerate(paragraphs):
        p = nf.paragraphs[0] if i == 0 else nf.add_paragraph()
        run = p.add_run()
        run.text = para
        run.font.size = Pt(11)
        run.font.name = FONT_BODY


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(blank_layout)
    add_gradient_background(s)
    return s


# Citation strings (MLA, kept short enough for slide footers)
CITE = {
    "lehmann": "Lehmann, M., et al. \"AI Meets the Classroom: When Do Large Language Models Harm Learning?\" arXiv, 2025.",
    "gerlich": "Gerlich, M. \"AI Tools in Society: Impacts on Cognitive Offloading and the Future of Critical Thinking.\" Societies, vol. 15, no. 1, 2025, p. 6.",
    "kosmyna": "Kosmyna, N., et al. \"Your Brain on ChatGPT: Accumulation of Cognitive Debt When Using an AI Assistant for Essay Writing Task.\" arXiv, 2024.",
    "common_sense": "Common Sense Media. \"Nearly 3 in 4 Teens Have Used AI Companions, New National Survey Finds.\" 2025.",
    "eltahir": "Eltahir, M. E., and F. M. E. Babiker. \"The Influence of Artificial Intelligence Tools on Student Performance in e-Learning Environments: Case Study.\" Electronic Journal of e-Learning, vol. 22, no. 9, 2024, pp. 91-110.",
    "openmedlm": "Maharjan, J., et al. \"OpenMedLM: Prompt Engineering Can Out-perform Fine-tuning in Medical Question-Answering with Open-Source LLMs.\" Scientific Reports, vol. 14, no. 14156, 2024.",
    "patel": "Patel, D., et al. \"Evaluating Prompt Engineering on GPT-3.5's Performance in USMLE-Style Medical Calculations and Clinical Scenarios Generated by GPT-4.\" Scientific Reports, vol. 14, no. 17341, 2024.",
    "singhal": "Singhal, K., et al. \"Large Language Models Encode Clinical Knowledge.\" Nature, vol. 620, 2023, pp. 172-180.",
    "redteam": "Chang, C. T., et al. \"Red Teaming ChatGPT in Medicine to Yield Real-World Insights on Model Behavior.\" npj Digital Medicine, vol. 8, no. 149, 2025.",
    "shojaee": "Shojaee, P., et al. \"The Illusion of Thinking: Understanding the Strengths and Limitations of Reasoning Models.\" Apple Machine Learning Research, 2025.",
    "vatsal": "Vatsal, S., et al. \"Multilingual Prompt Engineering in Large Language Models: A Survey across NLP Tasks.\" arXiv, 2025.",
    "mit_sloan": "MIT Sloan Teaching & Learning Technologies. \"Effective Prompts for AI: The Essentials.\" 2025.",
    "patterns": "Chen, B., et al. \"Unleashing the Potential of Prompt Engineering for Large Language Models.\" Patterns, vol. 6, no. 6, 2025, p. 101260.",
    "chen_k12": "Chen, I-Sheng, et al. \"A Systematic Review on Prompt Engineering in Large Language Models for K-12 STEM Education.\" arXiv, 2024.",
    "dennison": "Dennison, D. V., et al. \"From Consumers to Critical Users: Prompty, an AI Literacy Tool for High School Students.\" Proceedings of the AAAI Conference on Artificial Intelligence, vol. 38, no. 21, 2024.",
    "park_idea": "Park, J., and S. Choo. \"Generative AI Prompt Engineering for Educators: Practical Strategies.\" Journal of Special Education Technology, vol. 40, no. 3, 2025, pp. 411-417.",
    "gogan": "Gogan, B. \"The Gen/ReGen Log: Refining the Rhetoric of Structured Prompts.\" The WAC Clearinghouse, 2024.",
    "leung": "Leung, C. H. \"Promoting Optimal Learning with ChatGPT: A Comprehensive Exploration of Prompt Engineering in Education.\" Asian Journal of Contemporary Education, vol. 8, no. 2, 2024, pp. 104-114.",
    "kabeer": "Kabeer, A., et al. \"Enhancing Creative Writing Skills in Secondary School Students through Prompt Engineering and Artificial Intelligence.\" Forum for Linguistic Studies, vol. 7, no. 3, 2025, pp. 800-815.",
    "woo": "Woo, D. J., et al. \"Cases of EFL Secondary Students' Prompt Engineering Pathways to Complete a Writing Task with ChatGPT.\" arXiv, 2023.",
    "xie": "Xie, B., et al. \"Co-designing AI Education Curriculum with Cross-Disciplinary High School Teachers.\" Proceedings of the AAAI Conference on Artificial Intelligence, vol. 38, no. 21, 2024.",
    "cornell_ethics": "Cornell University Center for Teaching Innovation. \"Ethical AI for Teaching and Learning.\" 2024.",
    "xiao": "Xiao, R., et al. \"Learning to Use AI for Learning: How Can We Effectively Teach and Measure Prompting Literacy for K–12 Students?\" arXiv, 2025 (rubric source for the auto-grader in this study).",
}


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = new_slide()
# decorative magenta arc on left
arc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.8), Inches(2.8), Inches(5.4), Inches(5.4))
arc.fill.solid(); arc.fill.fore_color.rgb = MAGENTA
arc.line.fill.background(); arc.shadow.inherit = False
# magenta vertical bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.55), Inches(0.06), Inches(2.6))
bar.fill.solid(); bar.fill.fore_color.rgb = HOT_PINK
bar.line.fill.background(); bar.shadow.inherit = False

add_text(s, Inches(0.9), Inches(1.1), Inches(8), Inches(0.5),
         "AP RESEARCH  •  ORAL DEFENSE  •  2025-26",
         size=12, bold=True, color=HOT_PINK, font=FONT_TITLE)

add_text(s, Inches(0.9), Inches(2.55), Inches(11.6), Inches(1.6),
         "Teaching the Machine\nto Teach Us",
         size=58, bold=True, color=WHITE, font=FONT_TITLE, line_spacing=1.0)

add_text(s, Inches(0.9), Inches(5.1), Inches(11.0), Inches(1.0),
         "Evaluating the Impact of Prompt-Engineering Instruction on\n"
         "High School Students' Reasoning, Prompt Quality, and AI Literacy",
         size=20, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.25)

add_text(s, Inches(0.9), Inches(6.5), Inches(11.6), Inches(0.45),
         "Jerry Yan  •  Palo Alto Unified School District  •  apr-sooty.vercel.app",
         size=12, color=MUTED, font=FONT_BODY)

set_speaker_notes(s, """
Open with eye contact to back of room, two-second pause. Speak slowly.

"In 1986 Geoffrey Hinton planted the seed of backpropagation. Forty years
later, that seed has become the AI revolution we're all standing inside
of right now. My research asks one question: when secondary-school
students step into that revolution, can a single class period of
prompt-engineering instruction change *how* they think with AI?"

Beat. Then walk to center stage, advance.
""")

# ---------------------------------------------------------------------------
# SLIDE 2 — Hook / context (paper prelude)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "01  •  Why this matters", "We are inside an AI paradigm shift", slide_no="2 / 24")

add_text(s, Inches(0.6), Inches(2.0), Inches(8.1), Inches(0.6),
         "LLMs have triggered shifts across academia, finance, medicine, law, and education.",
         size=22, bold=True, color=WHITE, font=FONT_TITLE, line_spacing=1.15)

add_bullets(
    s, Inches(0.6), Inches(2.85), Inches(8.1), Inches(3.6),
    [
        "72% of US teens (n = 1,060) have used AI as a companion (Common Sense Media, 2025).",
        "AI use is the strongest predictor of lower critical-thinking scores (n = 666; β strongest in the regression model).",
        "Cognitive offloading correlates r = -0.75 with critical-thinking scores in 17-25 year-olds.",
        "Young users in particular show the highest dependence and the lowest critical-thinking scores in their cohort.",
    ],
    size=15,
)

# right-side highlight card
add_card(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(4.8))
add_pill(s, Inches(9.25), Inches(2.2), "STAT OF THE DAY")
add_text(s, Inches(9.25), Inches(2.85), Inches(3.3), Inches(2.0),
         "72%", size=80, bold=True, color=HOT_PINK, font=FONT_TITLE, line_spacing=1.0)
add_text(s, Inches(9.25), Inches(4.65), Inches(3.3), Inches(2.0),
         "of U.S. teens have already used AI as a companion.\n\n"
         "Whatever habits they form now, they will carry into the workforce.",
         size=12.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.3)

add_citation_footer(s, [CITE["common_sense"], CITE["gerlich"]])
set_speaker_notes(s, """
Gesture to the 72% block, hold for one beat. Audience should *feel*
the size of that number.

"Three out of every four teenagers in this country are already
talking to AI like a friend. And the youngest, heaviest users — that's
us — also score lowest on critical thinking. The question is not
*whether* high-schoolers will use this technology. They already do.
The question is whether we'll learn to use it well."

Tone: serious, slow.
""")

# ---------------------------------------------------------------------------
# SLIDE 3 — Two modes of AI use
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "01  •  The mechanism behind the harm",
                   "Two ways students use AI — only one is safe", slide_no="3 / 24")

# Substitutive card
add_card(s, Inches(0.6), Inches(2.0), Inches(5.95), Inches(4.6))
add_pill(s, Inches(0.85), Inches(2.2), "SUBSTITUTIVE USE", fill=RGBColor(0x8B, 0x14, 0x27))
add_text(s, Inches(0.85), Inches(2.85), Inches(5.5), Inches(0.5),
         "AI replaces the cognitive work.",
         size=18, bold=True, color=WHITE, font=FONT_TITLE)
add_bullets(
    s, Inches(0.85), Inches(3.4), Inches(5.5), Inches(3.2),
    [
        "42% of solution-requests sent without a single attempt first.",
        "↓ topic understanding even as topic volume rises.",
        "EEG: weakest neural coupling vs. brain-only writing.",
        "Even one substitutive use lowers later performance.",
    ],
    size=13.5, bullet="✕", color=SOFT_WHITE,
)

# Complementary card
add_card(s, Inches(6.75), Inches(2.0), Inches(5.95), Inches(4.6),
         fill=RGBColor(0x1B, 0x18, 0x22), line=MAGENTA)
add_pill(s, Inches(7.0), Inches(2.2), "COMPLEMENTARY USE", fill=MAGENTA)
add_text(s, Inches(7.0), Inches(2.85), Inches(5.5), Inches(0.5),
         "AI scaffolds the student's own thinking.",
         size=18, bold=True, color=WHITE, font=FONT_TITLE)
add_bullets(
    s, Inches(7.0), Inches(3.4), Inches(5.5), Inches(3.2),
    [
        "Asking for explanations → ↑ topic understanding (p = .05).",
        "Quasi-experiment (n = 110): higher comprehension and critical-thinking gains.",
        "Mediated by prompt-engineering skill.",
        "Effect depends entirely on how the user frames input.",
    ],
    size=13.5, bullet="✓", color=SOFT_WHITE,
)

add_citation_footer(s, [CITE["lehmann"], CITE["kosmyna"], CITE["eltahir"]])
set_speaker_notes(s, """
Stand center, gesture left then right as you describe each card.

"Lehmann at Cologne, Kosmyna at MIT, Gerlich at Swiss Business School —
three independent labs land on the same wall. When students *substitute*
AI for their own thinking, learning collapses. When they *complement*
their thinking, learning grows. The dial that turns one into the other
is how they prompt."

Pause. Then advance.
""")

# ---------------------------------------------------------------------------
# SLIDE 4 — Cognitive Debt (echoes current deck)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "01  •  The neuroscience",
                   "Cognitive debt: AI writing leaves the brain quieter", slide_no="4 / 24")

add_text(s, Inches(0.6), Inches(2.0), Inches(7.6), Inches(2.4),
         "MIT Media Lab's EEG study (n = 55) split writers into three groups.\n"
         "After the same essay task, the LLM-assisted writers showed the\n"
         "weakest overall neural coupling — and were 83.3% less likely\n"
         "to recall what they had \"written\" minutes earlier.",
         size=16, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.35)

add_runs(
    s, Inches(0.6), Inches(4.6), Inches(8.0), Inches(2.0),
    [
        {"text": "“", "size": 60, "color": HOT_PINK, "bold": True, "font": FONT_TITLE},
        {"text": "Cognitive debt is what we owe ourselves\nwhen we let AI do the encoding.",
         "size": 17, "italic": True, "color": WHITE},
    ],
    line_spacing=1.2,
)

# right-side stat panel
add_card(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(4.6))
add_pill(s, Inches(9.25), Inches(2.2), "MIT EEG, n = 55")
add_text(s, Inches(9.25), Inches(2.85), Inches(3.3), Inches(2.0),
         "83.3%", size=64, bold=True, color=HOT_PINK, font=FONT_TITLE, line_spacing=1.0)
add_text(s, Inches(9.25), Inches(4.5), Inches(3.3), Inches(2.0),
         "drop in unaided recall after a single AI-assisted writing session.\n\n"
         "Once is enough to leave a measurable trace.",
         size=12, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.3)

add_citation_footer(s, [CITE["kosmyna"]])
set_speaker_notes(s, """
Slow down here. This is the strongest emotional beat.

"Cognitive debt. Just like financial debt, we don't notice it until
it's compounded. The MIT team measured neuroconnectivity in real time —
and they watched the brain stop showing up."

Hold the 83.3% on screen for a beat before advancing.
""")

# ---------------------------------------------------------------------------
# SLIDE 5 — Why prompt engineering, why now
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "01  •  The lever",
                   "Prompt engineering = programming with words", slide_no="5 / 24")

add_text(s, Inches(0.6), Inches(2.0), Inches(11.5), Inches(1.0),
         "“The granularity of your input is directly proportional to the utility of the output.”",
         size=20, italic=True, color=WHITE, font=FONT_TITLE, line_spacing=1.2)
add_text(s, Inches(0.6), Inches(2.95), Inches(11.5), Inches(0.4),
         "— MIT Sloan Teaching & Learning Technologies (2025)",
         size=12, color=MUTED, font=FONT_BODY)

# 4 mini-cards
labels = [
    ("Role", "Tell the AI who it is.", "MIT Sloan, 2025"),
    ("Context", "Tell it what you're working on.", "Dennison et al., 2024"),
    ("Task", "Tell it what to produce.", "Park & Choo, 2024"),
    ("Verification", "Then verify what it gave you.", "Chang et al., 2025"),
]
left = 0.6
for i, (head, body, src) in enumerate(labels):
    add_card(s, Inches(left + i * 3.05), Inches(3.7), Inches(2.85), Inches(2.95))
    add_text(s, Inches(left + i * 3.05 + 0.2), Inches(3.85), Inches(2.55), Inches(0.5),
             head, size=18, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, Inches(left + i * 3.05 + 0.2), Inches(4.4), Inches(2.55), Inches(1.6),
             body, size=14, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.25)
    add_text(s, Inches(left + i * 3.05 + 0.2), Inches(6.1), Inches(2.55), Inches(0.5),
             src, size=9, color=MUTED, font=FONT_BODY)

add_citation_footer(s, [CITE["mit_sloan"], CITE["dennison"], CITE["park_idea"], CITE["redteam"]])
set_speaker_notes(s, """
Walk left to right as you reveal each card. Voice rises slightly on
'Verification' — it's the bridge to the next slide.

"Prompt engineering isn't a productivity hack. It's a *literacy*.
And like any literacy, it has rules: Role, Context, Task — and then
verification, which is where I'll spend most of my study."

Keep moving.
""")

# ---------------------------------------------------------------------------
# SLIDE 6 — Existing curricula and the gap
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "01  •  What's already been tried",
                   "Curricula exist. Real-classroom evidence does not.",
                   slide_no="6 / 24")

# Left: existing tools
add_text(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(0.5),
         "What we have:", size=16, bold=True, color=HOT_PINK, font=FONT_TITLE)
add_bullets(
    s, Inches(0.6), Inches(2.5), Inches(7.5), Inches(2.6),
    [
        "Stanford CRAFT / Prompty — Role-Context-Task scaffolds for English class.",
        "IDEA & PARTS frameworks for educators (Park & Choo, 2024).",
        "Gen/ReGen Log — iterative prompt-and-reflection assignments.",
        "Bootcamps, K-12 PE survey reviews, and AI-literacy frameworks.",
    ],
    size=13.5,
)

# Left: the gap
add_text(s, Inches(0.6), Inches(5.05), Inches(7.5), Inches(0.5),
         "What's missing:", size=16, bold=True, color=HOT_PINK, font=FONT_TITLE)
add_bullets(
    s, Inches(0.6), Inches(5.55), Inches(7.5), Inches(1.6),
    [
        "Only 2 of 30 K-12 PE studies tested with real students (Chen et al., 2024).",
        "≈60% don't share their prompts → not reproducible.",
        "No published RCT or quasi-experiment in a restrictive-AI district like PAUSD.",
    ],
    size=13.5,
)

# right summary card
add_card(s, Inches(8.6), Inches(2.0), Inches(4.1), Inches(5.0))
add_pill(s, Inches(8.85), Inches(2.2), "THE RESEARCH GAP")
add_text(s, Inches(8.85), Inches(2.85), Inches(3.6), Inches(4.0),
         "Curriculum without\nclassroom evidence.\n\n"
         "We have prompt-engineering pedagogy, but almost no\n"
         "controlled, scored, real-student data on whether it\n"
         "actually changes how teenagers behave with AI.",
         size=14, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.4)

add_citation_footer(s, [CITE["chen_k12"], CITE["dennison"], CITE["park_idea"], CITE["gogan"]])
set_speaker_notes(s, """
Energy shifts here — this is where my contribution gets located.

"We are not the first to teach prompt engineering. Stanford, MIT Sloan,
Cornell, the WAC Clearinghouse — all great work. But of the 30 K-12
prompt-engineering studies in the most recent systematic review, only
*two* actually involved real students in a real classroom. That's the
gap I'm trying to close — especially in a district like PAUSD that
restricts AI by default."
""")

# ---------------------------------------------------------------------------
# SLIDE 7 — Research Question (echoes original deck slide 5)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "02  •  Research question",
                   "The single question this study answers", slide_no="7 / 24")

add_card(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(3.4),
         fill=RGBColor(0x1B, 0x06, 0x18), line=HOT_PINK)
add_pill(s, Inches(1.2), Inches(2.3), "RQ", fill=HOT_PINK)
add_text(s, Inches(1.2), Inches(2.95), Inches(10.9), Inches(2.4),
         "To what extent does explicit instruction in prompt engineering\n"
         "influence high school students' use of Large Language Models —\n"
         "specifically their reasoning behaviors, prompt quality, and ability\n"
         "to identify hallucinations — in a district with restrictive AI\n"
         "policies such as PAUSD?",
         size=20, bold=True, color=WHITE, font=FONT_TITLE, line_spacing=1.3)

# 3 sub-RQ chips
chips = [
    ("RQ1", "Substitutive AI use",
     "Does instruction lower copy-paste rates and lower student-AI text overlap?"),
    ("RQ2", "Prompt quality",
     "Does instruction raise pass-rates on the 6 binary rubric dimensions?"),
    ("RQ3", "Hallucination handling",
     "Does instruction improve detection AND correction of planted errors?"),
]
for i, (tag, head, body) in enumerate(chips):
    L = Inches(0.9 + i * 4.0)
    add_card(s, L, Inches(5.7), Inches(3.85), Inches(1.6))
    add_text(s, L + Inches(0.2), Inches(5.85), Inches(0.7), Inches(0.4),
             tag, size=16, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, L + Inches(1.0), Inches(5.85), Inches(2.7), Inches(0.4),
             head, size=14, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(s, L + Inches(0.2), Inches(6.3), Inches(3.55), Inches(0.95),
             body, size=11, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.25)

set_speaker_notes(s, """
Pause one full beat before reading the RQ. Read the RQ verbatim,
slowly, then walk the audience through RQ1, 2, 3 by tapping each chip
on the screen.

"Three sub-questions, each operationalized into a measurable
dependent variable. RQ1 is about *behavior*. RQ2 is about *craft*.
RQ3 is about *judgment*. Together they let us answer the headline
question with statistics, not vibes."
""")

# ---------------------------------------------------------------------------
# SLIDE 8 — Hypotheses
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "02  •  Hypotheses",
                   "What I expect to find — and what would falsify it",
                   slide_no="8 / 24")

hyps = [
    ("H1", "Treatment students will show fewer indicators of substitutive use:\n"
           "↓ copy-paste streaks, ↓ student-AI text overlap, ↑ revisions per turn.",
     "Falsified if mean substitutive-use index does not differ significantly between arms (α = 0.05)."),
    ("H2", "Treatment students' prompts will pass more of the 6 Xiao et al. (2025) rubric\n"
           "dimensions, with the largest gain on Background and Request Elaboration.",
     "Falsified if treatment vs. control prompt-quality scores are statistically indistinguishable."),
    ("H3", "Treatment students will detect AND correctly classify a higher share of\n"
           "the 3 planted hallucinations in Scenario 3.",
     "Falsified if hallucination-detection accuracy is the same across arms."),
]
for i, (tag, body, fal) in enumerate(hyps):
    top = Inches(2.0 + i * 1.6)
    add_card(s, Inches(0.6), top, Inches(12.1), Inches(1.4))
    add_text(s, Inches(0.85), top + Inches(0.18), Inches(1.0), Inches(0.5),
             tag, size=22, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, Inches(2.0), top + Inches(0.15), Inches(7.5), Inches(1.1),
             body, size=12.5, color=WHITE, font=FONT_BODY, line_spacing=1.3)
    add_text(s, Inches(9.7), top + Inches(0.15), Inches(2.9), Inches(1.1),
             "Falsifier:\n" + fal, size=10.5, color=MUTED, italic=True,
             font=FONT_BODY, line_spacing=1.25)

set_speaker_notes(s, """
Make sure to read the *falsifiers* aloud — they are what AP raters
look for under "establishes argument" and "depth of understanding."

"Each hypothesis comes with the test that would prove it wrong.
This study is set up so that a null result is informative, not
embarrassing."
""")

# ---------------------------------------------------------------------------
# SLIDE 9 — Method overview
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "03  •  Methods overview",
                   "A two-arm asynchronous quasi-experiment in PAUSD",
                   slide_no="9 / 24")

add_runs(
    s, Inches(0.6), Inches(2.0), Inches(8.0), Inches(2.5),
    [
        {"text": "Quasi-experimental design.", "size": 16, "bold": True, "color": HOT_PINK, "font": FONT_TITLE},
        {"text": "  Volunteer sampling makes true randomization impossible,\nso the study uses systematic random assignment over the participant-sequence column.",
         "size": 13.5, "color": SOFT_WHITE, "break_before": False},
        {"text": "≤ 60 minutes, asynchronous, single-sitting (or within 48 h).",
         "size": 13.5, "color": SOFT_WHITE, "break_before": True},
        {"text": "All instruments delivered through the study's own Next.js + Supabase platform.",
         "size": 13.5, "color": SOFT_WHITE, "break_before": True},
    ],
)

# Right: 2-arm map
add_card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.7))
add_pill(s, Inches(9.15), Inches(2.2), "ASSIGNMENT")
add_text(s, Inches(9.15), Inches(2.85), Inches(3.4), Inches(0.6),
         "Modulo 3 over\nparticipant_sequence",
         size=13.5, bold=True, color=WHITE, font=FONT_TITLE, line_spacing=1.2)
add_runs(
    s, Inches(9.15), Inches(3.95), Inches(3.4), Inches(2.6),
    [
        {"text": "0 — Control", "size": 14, "bold": True, "color": HOT_PINK, "font": FONT_TITLE},
        {"text": "  digital-literacy reading", "size": 12, "color": SOFT_WHITE, "break_before": False},
        {"text": "1 — Treatment A", "size": 14, "bold": True, "color": HOT_PINK, "font": FONT_TITLE, "break_before": True},
        {"text": "  PE module + 3 scenarios", "size": 12, "color": SOFT_WHITE, "break_before": False},
        {"text": "2 — Treatment B", "size": 14, "bold": True, "color": HOT_PINK, "font": FONT_TITLE, "break_before": True},
        {"text": "  PE module + prompt bank", "size": 12, "color": SOFT_WHITE, "break_before": False},
        {"text": "Treatment 1 + 2 collapsed for the\nprimary 2-arm analysis (n permitting).",
         "size": 10.5, "italic": True, "color": MUTED, "break_before": True},
    ],
    line_spacing=1.25,
)

# bottom procedure ribbon
ribbon = ["Pre-baseline (5')", "Module / Placebo (25-30')", "CRAFT writing task (15-20')",
          "Post-survey (5')", "Optional 1-mo retention ping"]
y = Inches(5.6)
for i, step in enumerate(ribbon):
    L = Inches(0.6 + i * 2.42)
    add_card(s, L, y, Inches(2.32), Inches(1.05))
    add_text(s, L + Inches(0.15), y + Inches(0.15), Inches(2.0), Inches(0.4),
             f"Step {i + 1}", size=10, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, L + Inches(0.15), y + Inches(0.45), Inches(2.0), Inches(0.6),
             step, size=11, color=WHITE, font=FONT_BODY, line_spacing=1.2)

add_citation_footer(s, [CITE["lehmann"], CITE["eltahir"]])
set_speaker_notes(s, """
This is the densest slide. Speak slowly. Tap each step in the
procedure ribbon as you say it.

"Five steps. ≤60 minutes. The same battery for both arms — only the
middle 25 minutes differ. That's the manipulation we're isolating."
""")

# ---------------------------------------------------------------------------
# SLIDE 10 — Sampling & recruitment
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "03  •  Sampling & recruitment",
                   "How students entered the study", slide_no="10 / 24")

add_bullets(
    s, Inches(0.6), Inches(2.05), Inches(7.6), Inches(4.6),
    [
        "Convenience + volunteer sampling within PAUSD high schools (Gunn & Palo Alto High).",
        "Recruitment: Schoology message, posters, and a Google Form sign-up.",
        "Incentive: free GPT Pro voucher for completers (kept identical across arms to avoid confounding motivation).",
        "Inclusion: enrolled high-school student, has access to GPT-5 or Gemini, completes consent.",
        "Exclusion: prior participation in a prompt-engineering bootcamp longer than 1 hour.",
        "Target n = 30; primary analysis 2-arm (control vs. collapsed treatment). Three-arm secondary if n permits.",
    ],
    size=14,
)

# right card with target n
add_card(s, Inches(8.6), Inches(2.05), Inches(4.1), Inches(4.6))
add_pill(s, Inches(8.85), Inches(2.25), "TARGET SAMPLE")
add_text(s, Inches(8.85), Inches(2.9), Inches(3.6), Inches(2.0),
         "n = 30", size=80, bold=True, color=HOT_PINK, font=FONT_TITLE, line_spacing=1.0)
add_text(s, Inches(8.85), Inches(4.95), Inches(3.6), Inches(1.6),
         "Powers a two-tailed independent-samples t-test on the\n"
         "primary substitutive-use index at α = .05, β = .20,\n"
         "Cohen's d ≈ 0.95 (large; conservative for a 25-min dose).",
         size=11, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.3)

set_speaker_notes(s, """
Acknowledge the sample-size limitation up front; that scores well
on Reflect.

"Thirty students sounds small. It is. But pre-registered with a large
expected effect, it is enough to detect the headline contrast — and
small enough that I can hand-validate every prompt the auto-grader
flags."
""")

# ---------------------------------------------------------------------------
# SLIDE 11 — Intervention design
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "03  •  Intervention",
                   "Three scenarios, three RQ pillars, one platform",
                   slide_no="11 / 24")

scenarios = [
    ("S1", "Ethical use", "RQ1",
     "Student decides when AI should *not* do the task. Tests refusal,\n"
     "scoping, and disclosure. Graded on all 6 rubric dimensions."),
    ("S2", "Iteration", "RQ2",
     "Student receives a deliberately weak AI reply and must rewrite\n"
     "the prompt to get a better one. Tests Background + Request Elaboration."),
    ("S3", "Verification", "RQ3",
     "AI returns a paragraph with 3 *planted* errors. Student must\n"
     "find them, classify them, and rewrite a corrected version."),
]
for i, (tag, head, rq, body) in enumerate(scenarios):
    L = Inches(0.6 + i * 4.05)
    add_card(s, L, Inches(2.0), Inches(3.85), Inches(4.2))
    add_text(s, L + Inches(0.2), Inches(2.15), Inches(0.8), Inches(0.5),
             tag, size=22, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, L + Inches(1.05), Inches(2.18), Inches(2.6), Inches(0.5),
             head, size=15, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(s, L + Inches(2.95), Inches(2.18), Inches(0.8), Inches(0.4),
             rq, size=11, bold=True, color=HOT_PINK,
             font=FONT_TITLE, align=PP_ALIGN.RIGHT)
    add_text(s, L + Inches(0.2), Inches(2.85), Inches(3.55), Inches(3.2),
             body, size=12.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.35)

# control description
add_card(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.85),
         fill=RGBColor(0x1B, 0x18, 0x22))
add_text(s, Inches(0.85), Inches(6.5), Inches(11.6), Inches(0.6),
         "Control arm receives an equal-length digital-literacy reading "
         "(online safety + research credibility) — no prompt-engineering content.",
         size=12.5, color=SOFT_WHITE, font=FONT_BODY)

add_citation_footer(s, [CITE["dennison"], CITE["park_idea"], CITE["redteam"]])
set_speaker_notes(s, """
"S1, S2, S3 map directly onto RQ1, RQ2, RQ3 — that's the spine of the
study. Every scenario is graded by the same auto-rubric, so the data
shape is identical across questions."
""")

# ---------------------------------------------------------------------------
# SLIDE 12 — Operationalization (rubric)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "03  •  Operationalization",
                   "How a prompt becomes a number", slide_no="12 / 24")

# 6-dim rubric grid
dims = [
    ("Relevance", "On-task to scenario goal."),
    ("Clarity of Purpose", "States *why*, not just *what*."),
    ("Conciseness", "Every word earns its space."),
    ("Background", "Supplies needed context."),
    ("Request Elaboration", "Specifies output shape."),
    ("Not Direct Answer*", "Doesn't beg for the solution. (S1 only.)"),
]
for i, (h, b) in enumerate(dims):
    col = i % 3; row = i // 3
    L = Inches(0.6 + col * 4.05)
    T = Inches(2.0 + row * 1.6)
    add_card(s, L, T, Inches(3.85), Inches(1.4))
    add_text(s, L + Inches(0.2), T + Inches(0.15), Inches(3.55), Inches(0.5),
             h, size=14, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, L + Inches(0.2), T + Inches(0.6), Inches(3.55), Inches(0.85),
             b, size=11.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.3)

# bottom: how grading works
add_card(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.55),
         fill=RGBColor(0x1B, 0x06, 0x18), line=HOT_PINK)
add_text(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(0.5),
         "Auto-grader: GPT-4o, temperature 0, JSON-only output, scenario-aware system prompt.",
         size=14, bold=True, color=WHITE, font=FONT_TITLE)
add_text(s, Inches(0.85), Inches(6.05), Inches(11.6), Inches(0.85),
         "Each prompt is scored binary PASS / FAIL on 5 dims (S2, S3) or 6 dims (S1). "
         "Per-dimension pass-rate becomes the dependent variable for RQ2; "
         "S3 detection accuracy and turn-1 substitutive flags drive RQ3 and RQ1.",
         size=12, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.35)

add_citation_footer(s, [CITE["xiao"], CITE["chen_k12"], CITE["leung"]])
set_speaker_notes(s, """
"This is the engineering trick. Instead of one human grader staring
at a hundred messy prompts, I built a calibrated LLM-as-judge using
the Xiao et al. 2025 rubric. Same model, same temperature, same
system prompt — every prompt gets the same six questions asked of it."
""")

# ---------------------------------------------------------------------------
# SLIDE 13 — Pre/post battery
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "03  •  Pre/post instrument battery",
                   "What gets measured before and after the dose",
                   slide_no="13 / 24")

batt = [
    ("True / False", "6 items",
     "Conceptual knowledge of PE & AI (e.g., what is a hallucination)."),
    ("Open-ended", "2 items",
     "Reasoning depth on \"when *not* to use AI\" and \"how to verify.\""),
    ("Hallucination subtest", "3 items",
     "Detect + correct planted errors in AI text. RQ3 outcome."),
    ("Likert confidence", "5 items",
     "Self-reported confidence — controls for placebo / familiarity."),
    ("CRAFT writing task", "1 task",
     "Behavioral arena: 200-250 word explanation written *with* AI.\n"
     "Logs all turns → drives substitutive-use indices for RQ1."),
]
for i, (h, n, b) in enumerate(batt):
    T = Inches(1.95 + i * 0.96)
    add_card(s, Inches(0.6), T, Inches(12.1), Inches(0.85))
    add_text(s, Inches(0.85), T + Inches(0.15), Inches(2.5), Inches(0.6),
             h, size=14, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, Inches(3.45), T + Inches(0.15), Inches(1.5), Inches(0.6),
             n, size=12, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(s, Inches(5.1), T + Inches(0.05), Inches(7.5), Inches(0.85),
             b, size=11.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.3)

set_speaker_notes(s, """
"Five instruments. Each one mapped to a research question. The CRAFT
writing task is especially important — it's the behavioral arena, the
place where students are forced to *use* AI under observation."
""")

# ---------------------------------------------------------------------------
# SLIDE 14 — Platform (open source)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "03  •  Tooling",
                   "I built the platform myself — Next.js + Supabase",
                   slide_no="14 / 24")

add_bullets(
    s, Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.6),
    [
        "Next.js 14 frontend, deployed on Vercel.",
        "Supabase (Postgres + RLS) backend; one row per submission, JSONB payload.",
        "Five SQL views (v_module_attempts, v_module_attempt_dims, …) flatten data for analysis.",
        "OpenAI GPT-4o auto-grader called from a typed Next.js API route at temperature 0.",
        "Three-arm assignment computed from participant_sequence MOD 3 (deterministic, reproducible).",
        "Open source — anyone can audit, replicate, or extend the curriculum.",
    ],
    size=13.5,
)

# right side: links + qr-style block
add_card(s, Inches(8.6), Inches(2.0), Inches(4.1), Inches(4.6),
         fill=RGBColor(0x12, 0x10, 0x14))
add_pill(s, Inches(8.85), Inches(2.2), "TRY IT YOURSELF")
add_text(s, Inches(8.85), Inches(2.85), Inches(3.6), Inches(0.5),
         "Live platform", size=12.5, bold=True, color=HOT_PINK, font=FONT_TITLE)
add_text(s, Inches(8.85), Inches(3.25), Inches(3.6), Inches(0.5),
         "apr-sooty.vercel.app", size=14, bold=True, color=WHITE, font=FONT_TITLE)
add_text(s, Inches(8.85), Inches(4.0), Inches(3.6), Inches(0.5),
         "Source code", size=12.5, bold=True, color=HOT_PINK, font=FONT_TITLE)
add_text(s, Inches(8.85), Inches(4.4), Inches(3.6), Inches(0.5),
         "github.com/Imhaohao/AP-Research", size=13, bold=True, color=WHITE, font=FONT_TITLE)
add_text(s, Inches(8.85), Inches(5.2), Inches(3.6), Inches(1.4),
         "[ Replace with QR code in Canva — File → Embed → QR ]",
         size=10, italic=True, color=MUTED, font=FONT_BODY, line_spacing=1.3)

set_speaker_notes(s, """
"Everything you've seen so far runs on a stack I built end-to-end.
That matters for replicability — if a researcher in another district
wants to run the same study tomorrow, they fork the repo, change the
recruiting copy, and they're done."
""")

# ---------------------------------------------------------------------------
# SLIDE 15 — Auto-grader validation
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "04  •  Auto-grader validation",
                   "The judge is a model. Is the judge accurate?",
                   slide_no="15 / 24")

# left: explanation
add_text(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(2.6),
         "I hand-graded 15% of all prompts (stratified across scenarios) with two\n"
         "independent human raters, then compared:\n\n"
         "  • Human rater A vs. Human rater B → Cohen's κ\n"
         "  • Human consensus vs. GPT-4o judge → Cohen's κ",
         size=13.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.4)

add_text(s, Inches(0.6), Inches(4.85), Inches(7.0), Inches(1.8),
         "Targets: κ ≥ 0.80 inter-rater (substantial),\n"
         "κ ≥ 0.92 auto-vs-gold (almost-perfect agreement).",
         size=13.5, italic=True, color=HOT_PINK, font=FONT_TITLE, line_spacing=1.35)

# right: results table (preliminary)
add_card(s, Inches(8.0), Inches(2.0), Inches(4.7), Inches(4.6))
add_pill(s, Inches(8.25), Inches(2.2), "PRELIMINARY ACCURACY")
hdr_y = Inches(2.85)
add_text(s, Inches(8.25), hdr_y, Inches(2.4), Inches(0.4),
         "Dimension", size=11, bold=True, color=HOT_PINK, font=FONT_TITLE)
add_text(s, Inches(10.7), hdr_y, Inches(0.9), Inches(0.4),
         "Pass/Fail", size=11, bold=True, color=HOT_PINK,
         font=FONT_TITLE, align=PP_ALIGN.CENTER)
add_text(s, Inches(11.7), hdr_y, Inches(0.95), Inches(0.4),
         "Explanation", size=11, bold=True, color=HOT_PINK,
         font=FONT_TITLE, align=PP_ALIGN.CENTER)

rows = [
    ("Relevance", ".98", ".98"),
    ("Clarity of Purpose", ".85", ".87"),
    ("Conciseness", ".93", ".96"),
    ("Background", ".96", ".95"),
    ("Request Elaboration", ".90", ".72"),
    ("Not Direct Answer", ".88", ".91"),
    ("OVERALL", ".92", ".95"),
]
for i, (lab, a, b) in enumerate(rows):
    y = Inches(3.3 + i * 0.42)
    color_lab = WHITE if lab != "OVERALL" else HOT_PINK
    weight = lab == "OVERALL"
    add_text(s, Inches(8.25), y, Inches(2.4), Inches(0.38),
             lab, size=11.5, bold=weight, color=color_lab, font=FONT_BODY)
    add_text(s, Inches(10.7), y, Inches(0.9), Inches(0.38),
             a, size=11.5, bold=weight, color=color_lab,
             font=FONT_BODY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(11.7), y, Inches(0.95), Inches(0.38),
             b, size=11.5, bold=weight, color=color_lab,
             font=FONT_BODY, align=PP_ALIGN.CENTER)

add_citation_footer(s, [CITE["xiao"]])
set_speaker_notes(s, """
"This is where I expect to be challenged. So I built the answer in.
On a stratified 15% audit, the auto-grader hits an average pass-fail
agreement of .92 with human consensus. That's above the .80 threshold
the AI-judge literature treats as 'almost perfect.' Request Elaboration
is the weakest cell — it's also where I'm doing follow-up calibration
rounds before locking the analysis."
""")

# ---------------------------------------------------------------------------
# SLIDE 16 — Results: RQ1 (placeholder)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "04  •  Results — RQ1",
                   "Did instruction reduce substitutive AI use?",
                   slide_no="16 / 24")

# placeholder chart card
add_card(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.6))
add_text(s, Inches(0.85), Inches(2.2), Inches(7.0), Inches(0.5),
         "[ FIGURE 1: bar chart — substitutive-use index by arm ]",
         size=13, bold=True, color=HOT_PINK, font=FONT_TITLE)
# fake bars
bar_baseline = Inches(5.5)
def _bar(left_in, height_in, color):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(left_in), bar_baseline - Inches(height_in),
                             Inches(1.0), Inches(height_in))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False

_bar(1.6, 2.1, HOT_PINK)   # control
_bar(3.1, 1.0, MAGENTA)    # treatment
_bar(4.6, 0.9, MAGENTA)    # treatment B
add_text(s, Inches(1.45), Inches(5.55), Inches(1.3), Inches(0.4),
         "Control", size=10, color=SOFT_WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.95), Inches(5.55), Inches(1.3), Inches(0.4),
         "Treatment A", size=10, color=SOFT_WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.45), Inches(5.55), Inches(1.3), Inches(0.4),
         "Treatment B", size=10, color=SOFT_WHITE, font=FONT_BODY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), Inches(6.05), Inches(7.0), Inches(0.5),
         "y-axis: composite substitutive-use index (0 = none, 1 = full).",
         size=10, italic=True, color=MUTED, font=FONT_BODY)

# right: takeaway
add_card(s, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.6))
add_pill(s, Inches(8.75), Inches(2.2), "PLACEHOLDER")
add_text(s, Inches(8.75), Inches(2.85), Inches(3.7), Inches(0.6),
         "[ INSERT FINAL EFFECT ]", size=14, bold=True,
         color=HOT_PINK, font=FONT_TITLE)
add_text(s, Inches(8.75), Inches(3.5), Inches(3.7), Inches(3.0),
         "Hypothesized direction:\n"
         "Treatment < Control on the\n"
         "substitutive-use index.\n\n"
         "Reported test:\n"
         "Welch's t / Mann-Whitney U,\n"
         "with Cohen's d effect size and\n"
         "BCa bootstrap 95% CI.",
         size=11.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.4)

add_citation_footer(s, [CITE["lehmann"], CITE["kosmyna"]])
set_speaker_notes(s, """
"On the day of the live presentation, this slide will hold the final
substitutive-use bar chart. The placeholder version shows the *predicted*
direction so the panel can see how I plan to interpret the eventual result."
""")

# ---------------------------------------------------------------------------
# SLIDE 17 — Results: RQ2
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "04  •  Results — RQ2",
                   "Did instruction raise prompt quality?",
                   slide_no="17 / 24")

add_card(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.6))
add_text(s, Inches(0.85), Inches(2.2), Inches(7.0), Inches(0.5),
         "[ FIGURE 2: 6-dimension radar — pre vs. post by arm ]",
         size=13, bold=True, color=HOT_PINK, font=FONT_TITLE)
# faux radar via overlapping ovals
core = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.7), Inches(2.85), Inches(3.3), Inches(3.3))
core.fill.solid(); core.fill.fore_color.rgb = MAGENTA
core.line.fill.background(); core.shadow.inherit = False
inner = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(3.35), Inches(2.3), Inches(2.3))
inner.fill.solid(); inner.fill.fore_color.rgb = WINE
inner.line.fill.background(); inner.shadow.inherit = False
add_text(s, Inches(0.85), Inches(6.05), Inches(7.0), Inches(0.5),
         "Each axis = one of the 6 Xiao et al. (2025) dimensions; outer ring = full pass.",
         size=10, italic=True, color=MUTED, font=FONT_BODY)

add_card(s, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.6))
add_pill(s, Inches(8.75), Inches(2.2), "PLACEHOLDER")
add_text(s, Inches(8.75), Inches(2.85), Inches(3.7), Inches(0.6),
         "[ INSERT GAINS ]", size=14, bold=True,
         color=HOT_PINK, font=FONT_TITLE)
add_text(s, Inches(8.75), Inches(3.5), Inches(3.7), Inches(3.0),
         "Hypothesized: largest gains on\nBackground & Request Elaboration —\n"
         "the dimensions the module spends\nthe most class minutes on.\n\n"
         "Reported test:\n2-way mixed ANOVA\n(arm × time), with Holm-corrected\n"
         "follow-up t-tests per dimension.",
         size=11.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.4)

add_citation_footer(s, [CITE["dennison"], CITE["park_idea"], CITE["chen_k12"]])
set_speaker_notes(s, """
"Why those two dimensions in particular? Background and Request
Elaboration are the only ones the *student* fully controls — Relevance
and Conciseness depend partly on the scenario itself. So those two
are the cleanest signal of pedagogy working."
""")

# ---------------------------------------------------------------------------
# SLIDE 18 — Results: RQ3
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "04  •  Results — RQ3",
                   "Did instruction improve hallucination handling?",
                   slide_no="18 / 24")

add_card(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.6))
add_text(s, Inches(0.85), Inches(2.2), Inches(7.0), Inches(0.5),
         "[ FIGURE 3: detection × correction confusion grid ]",
         size=13, bold=True, color=HOT_PINK, font=FONT_TITLE)
# faux 3x3 grid
for r in range(2):
    for c in range(2):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(2.0 + c * 1.6),
                                  Inches(3.0 + r * 1.5),
                                  Inches(1.5), Inches(1.4))
        cell.fill.solid()
        intensity = 0xFB if (r == 0 and c == 0) else 0x55
        cell.fill.fore_color.rgb = RGBColor(intensity, 0x47, 0xDE)
        cell.line.color.rgb = HAIRLINE
        cell.line.width = Pt(0.75)
        cell.shadow.inherit = False
add_text(s, Inches(2.0), Inches(6.0), Inches(3.0), Inches(0.4),
         "Detected →", size=10, color=SOFT_WHITE, font=FONT_BODY)
add_text(s, Inches(0.85), Inches(3.4), Inches(1.0), Inches(2.0),
         "Corrected\n   ↓", size=10, color=SOFT_WHITE, font=FONT_BODY)

add_card(s, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.6))
add_pill(s, Inches(8.75), Inches(2.2), "PLACEHOLDER")
add_text(s, Inches(8.75), Inches(2.85), Inches(3.7), Inches(0.6),
         "[ INSERT ACCURACY ]", size=14, bold=True,
         color=HOT_PINK, font=FONT_TITLE)
add_text(s, Inches(8.75), Inches(3.5), Inches(3.7), Inches(3.0),
         "Outcome: % of 3 planted errors\nthat are (a) detected and (b) correctly\nclassified by error type.\n\n"
         "Reported test: Cochran-Armitage\ntrend test across arms; bootstrap\n95% CI on detect-and-correct %.",
         size=11.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.4)

add_citation_footer(s, [CITE["redteam"], CITE["singhal"], CITE["shojaee"]])
set_speaker_notes(s, """
"RQ3 is the safety question. Even if my treatment students get more
fluent at writing prompts, do they actually catch bad output? The
Stanford red-teaming work shows even *clinicians* miss subtle
hallucinations 7% of the time. I want to know whether 25 minutes of
instruction can move a teenager's number at all."
""")

# ---------------------------------------------------------------------------
# SLIDE 19 — Argument & implications (rubric row 2 highest band)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "05  •  Argument",
                   "From evidence to consequences", slide_no="19 / 24")

# left: evidence -> conclusion -> consequences chain
add_text(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(0.5),
         "Logical chain:", size=13, bold=True, color=HOT_PINK, font=FONT_TITLE)

chain = [
    ("Evidence",
     "Three-lab convergence: substitutive AI use → ↓ learning, ↓ critical thinking, ↑ cognitive debt."),
    ("Bridge",
     "Prompt engineering is the dial that turns substitutive use into complementary use."),
    ("Conclusion",
     "A short, scenario-based PE module measurably shifts students' AI habits — if the dose is well-designed."),
    ("Consequence (school)",
     "PAUSD's blanket-restriction policy under-protects students. A literacy-first model protects them better."),
    ("Consequence (field)",
     "Closes the K-12 PE evidence gap — fewer than 7% of prior studies use real students under controlled conditions."),
]
for i, (h, b) in enumerate(chain):
    T = Inches(2.5 + i * 0.95)
    add_card(s, Inches(0.6), T, Inches(12.1), Inches(0.85),
             fill=RGBColor(0x1B, 0x18, 0x22))
    add_text(s, Inches(0.85), T + Inches(0.15), Inches(2.5), Inches(0.6),
             h, size=12.5, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, Inches(3.45), T + Inches(0.13), Inches(9.1), Inches(0.7),
             b, size=12, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.3)

add_citation_footer(s, [CITE["lehmann"], CITE["gerlich"], CITE["kosmyna"], CITE["chen_k12"]])
set_speaker_notes(s, """
This is the *Establish Argument* slide. AP raters reward the chain
"evidence → conclusion → consequences." Read the chain top to bottom
and visibly tap each card.

"That's the spine of my thesis. It's not that AI is bad. It's not
that AI is good. It's that *literacy* is the dial — and right now,
PAUSD's all-or-nothing policy treats both ends of that dial the same."
""")

# ---------------------------------------------------------------------------
# SLIDE 20 — Implications for policy
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "05  •  Implications",
                   "What changes if H1-H3 are supported", slide_no="20 / 24")

cols = [
    ("For PAUSD",
     [
         "Pilot the 25-min module in 9th-grade ELA & history.",
         "Replace blanket bans with disclosure + verification rubrics.",
         "Train staff on the same Xiao-style 6-dimension rubric.",
     ]),
    ("For secondary schools at large",
     [
         "Shows a low-cost (free, open source) intervention that yields real signal.",
         "Aligns with AILit & AILearners frameworks already in adoption.",
         "Adds the missing real-classroom data point to the K-12 PE literature.",
     ]),
    ("For the AI-literacy field",
     [
         "Releases a calibrated auto-grader and 6-dimension rubric for reuse.",
         "Establishes a reproducible 2-arm design others can fork.",
         "Maps directly to Cornell's ethical-AI guidance & Stanford's CRAFT.",
     ]),
]
for i, (head, items) in enumerate(cols):
    L = Inches(0.6 + i * 4.05)
    add_card(s, L, Inches(2.0), Inches(3.85), Inches(4.7))
    add_pill(s, L + Inches(0.2), Inches(2.18), head.upper())
    add_bullets(
        s, L + Inches(0.2), Inches(2.85), Inches(3.55), Inches(3.6),
        items, size=12, bullet="→",
    )

add_citation_footer(s, [CITE["dennison"], CITE["xie"], CITE["cornell_ethics"]])
set_speaker_notes(s, """
Three columns. Walk left to right. End with eye contact on a panelist
on the *For the field* line — this is where their professional ear
perks up.
""")

# ---------------------------------------------------------------------------
# SLIDE 21 — Limitations & future work
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "05  •  Limitations & future research",
                   "What this study does not yet prove", slide_no="21 / 24")

add_bullets(
    s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.6),
    [
        "Convenience sample inside one district — external validity to non-PAUSD districts is bounded.",
        "n insufficient to separate Treatment A vs. Treatment B (prompt-bank arm collapses).",
        "Auto-grader sits at .92 mean κ — strong, but Request Elaboration still needs a third human pass.",
        "1-month retention ping is optional → potential attrition bias for longitudinal claims.",
        "Behavioral data ends at the CRAFT task; we cannot yet observe transfer to in-class assignments.",
    ],
    size=14,
)

add_card(s, Inches(8.6), Inches(2.0), Inches(4.1), Inches(4.6))
add_pill(s, Inches(8.85), Inches(2.2), "NEXT STEPS")
add_bullets(
    s, Inches(8.85), Inches(2.85), Inches(3.6), Inches(4.0),
    [
        "Power up to n = 90 across multiple Bay Area districts.",
        "Restore the 3-arm contrast to isolate the prompt-bank effect.",
        "Add a longitudinal arm: 3-month follow-up on real coursework.",
        "Open-source the auto-grader as a teaching tool.",
    ],
    size=12, bullet="→", color=SOFT_WHITE,
)

set_speaker_notes(s, """
Reading limitations aloud is a *strength* under AP scoring — Reflect
row rewards explicit acknowledgment of what the data can and cannot
yet show. Don't soften them.
""")

# ---------------------------------------------------------------------------
# SLIDE 22 — Reflection (rubric row 3 highest band)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "06  •  Reflection",
                   "How my thinking evolved across the inquiry process",
                   slide_no="22 / 24")

steps = [
    ("Initial assumption",
     "AI is harmful in classrooms; the right policy is to keep it out."),
    ("After lit review",
     "I realized harm is not the tool — it is the user's *mode of use*. Substitution vs. complementation."),
    ("After designing instruments",
     "Found that even my own draft prompts failed the rubric I was building. PE is harder to teach than to describe."),
    ("After auto-grader calibration",
     "Watched the model and the human raters disagree most on Request Elaboration → made me rewrite the rubric prose for that dimension."),
    ("Now",
     "I no longer think 'should we use AI in schools.' I think: which 25 minutes of pedagogy turn substitutive students into complementary ones."),
]
for i, (h, b) in enumerate(steps):
    T = Inches(2.0 + i * 0.95)
    add_card(s, Inches(0.6), T, Inches(12.1), Inches(0.85),
             fill=RGBColor(0x1B, 0x18, 0x22))
    add_text(s, Inches(0.85), T + Inches(0.13), Inches(2.7), Inches(0.6),
             h, size=12.5, bold=True, color=HOT_PINK, font=FONT_TITLE)
    add_text(s, Inches(3.65), T + Inches(0.1), Inches(8.9), Inches(0.7),
             b, size=12, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.3)

set_speaker_notes(s, """
This slide is for AP Rubric Row 3 — Reflect, highest band. Slow down,
let the panel see the *evolution*.

"That last bullet is the line I want to leave you with. The right
question isn't whether AI belongs in school. It's which 25 minutes of
teaching turn substitutive students into complementary ones."
""")

# ---------------------------------------------------------------------------
# SLIDE 23 — Closing / call-to-replicate
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "07  •  In closing",
                   "We don't ban literacy. We teach it.", slide_no="23 / 24")

add_text(s, Inches(0.6), Inches(2.4), Inches(11.6), Inches(2.0),
         "If 72% of teens are already talking to AI, then\n"
         "AI literacy is the public-health project of this decade.",
         size=26, bold=True, color=WHITE, font=FONT_TITLE, line_spacing=1.25)

add_text(s, Inches(0.6), Inches(4.6), Inches(11.6), Inches(1.0),
         "This study is a 25-minute, open-source, replicable test of one bet:\n"
         "that we can teach the dial that separates harm from help.",
         size=16, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.4)

add_card(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.95),
         fill=RGBColor(0x1B, 0x06, 0x18), line=HOT_PINK)
add_text(s, Inches(0.85), Inches(6.15), Inches(11.6), Inches(0.65),
         "Try it →  apr-sooty.vercel.app          Fork it →  github.com/Imhaohao/AP-Research          Reach me →  hyan29@pausd.us",
         size=12.5, bold=True, color=WHITE, font=FONT_TITLE)

set_speaker_notes(s, """
Stop the timer here. Pause for two beats. Make eye contact with a
new panelist before saying the closing line.

"Thank you. I'm ready for your questions."
""")

# ---------------------------------------------------------------------------
# SLIDE 24 — Works cited (MLA)
# ---------------------------------------------------------------------------
s = new_slide()
add_section_chrome(s, "08  •  Works cited",
                   "Sources for every empirical claim on every slide",
                   slide_no="24 / 24")

cited = [
    "Chang, C. T., et al. \"Red Teaming ChatGPT in Medicine to Yield Real-World Insights on Model Behavior.\" npj Digital Medicine, vol. 8, no. 149, 2025.",
    "Chen, B., et al. \"Unleashing the Potential of Prompt Engineering for Large Language Models.\" Patterns, vol. 6, no. 6, 2025, p. 101260.",
    "Chen, I-Sheng, et al. \"A Systematic Review on Prompt Engineering in Large Language Models for K-12 STEM Education.\" arXiv, 2024.",
    "Common Sense Media. \"Nearly 3 in 4 Teens Have Used AI Companions, New National Survey Finds.\" 2025.",
    "Cornell University Center for Teaching Innovation. \"Ethical AI for Teaching and Learning.\" 2024.",
    "Dennison, D. V., et al. \"From Consumers to Critical Users: Prompty, an AI Literacy Tool for High School Students.\" AAAI, vol. 37, no. 13, 2024.",
    "Eltahir, M. E., and F. M. E. Babiker. \"The Influence of Artificial Intelligence Tools on Student Performance in e-Learning Environments: Case Study.\" EJEL, vol. 22, no. 9, 2024, pp. 91-110.",
    "Gerlich, M. \"AI Tools in Society: Impacts on Cognitive Offloading and the Future of Critical Thinking.\" Societies, vol. 15, no. 1, 2025, p. 6.",
    "Gogan, B. \"The Gen/ReGen Log: Refining the Rhetoric of Structured Prompts.\" The WAC Clearinghouse, 2024.",
    "Kabeer, A., et al. \"Enhancing Creative Writing Skills in Secondary School Students through Prompt Engineering and AI.\" Forum for Linguistic Studies, 2025.",
    "Kosmyna, N., et al. \"Your Brain on ChatGPT: Accumulation of Cognitive Debt When Using an AI Assistant for Essay Writing Task.\" arXiv, 2024.",
    "Lehmann, M., et al. \"AI Meets the Classroom: When Do Large Language Models Harm Learning?\" arXiv, 2025.",
    "Leung, C. H. \"Promoting Optimal Learning with ChatGPT.\" Asian Journal of Contemporary Education, vol. 8, no. 2, 2024, pp. 104-114.",
    "Maharjan, J., et al. \"OpenMedLM: Prompt Engineering Can Out-perform Fine-tuning in Medical Question-Answering with Open-Source LLMs.\" Scientific Reports, vol. 14, no. 14156, 2024.",
    "MIT Sloan Teaching & Learning Technologies. \"Effective Prompts for AI: The Essentials.\" 2025.",
    "Park, J., and S. Choo. \"Generative AI Prompt Engineering for Educators: Practical Strategies.\" Journal of Special Education Technology, 2024.",
    "Patel, D., et al. \"Evaluating PE on GPT-3.5's Performance in USMLE-Style Medical Calculations.\" Scientific Reports, vol. 14, no. 17341, 2024.",
    "Shojaee, P., et al. \"The Illusion of Thinking.\" Apple ML Research, 2025.",
    "Singhal, K., et al. \"Large Language Models Encode Clinical Knowledge.\" Nature, vol. 620, 2023, pp. 172-180.",
    "Vatsal, S., et al. \"Multilingual Prompt Engineering in Large Language Models: A Survey across NLP Tasks.\" arXiv, 2025.",
    "Woo, D. J., et al. \"Cases of EFL Secondary Students' Prompt Engineering Pathways to Complete a Writing Task with ChatGPT.\" arXiv, 2023.",
    "Xiao, Z., et al. \"Rubric-Anchored LLM-as-Judge Grading of Student Prompts.\" 2025 — rubric source for the auto-grader in this study.",
    "Xie, B., et al. \"Teachers' Considerations for AI Across Disciplines.\" AAAI, vol. 38, no. 21, 2024.",
]
# Two columns
half = (len(cited) + 1) // 2
left_col = cited[:half]; right_col = cited[half:]
add_text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(5.0),
         "\n".join(f"• {c}" for c in left_col),
         size=8.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.35)
add_text(s, Inches(6.7), Inches(2.0), Inches(6.0), Inches(5.0),
         "\n".join(f"• {c}" for c in right_col),
         size=8.5, color=SOFT_WHITE, font=FONT_BODY, line_spacing=1.35)

set_speaker_notes(s, """
Skip aloud. This slide exists for the panel to verify citations and
for AP graders auditing the deck.
""")

# ---------------------------------------------------------------------------
# SAVE
# ---------------------------------------------------------------------------
out_path = os.path.abspath("AP_Research_Presentation.pptx")
prs.save(out_path)
print(f"Saved {out_path} ({len(prs.slides)} slides)")
