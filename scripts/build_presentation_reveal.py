from __future__ import annotations

import os
from pathlib import Path
from typing import Iterable, Sequence

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# 16:9 Canva-friendly
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WINE = RGBColor(0x34, 0x01, 0x1F)
BLACK = RGBColor(0x00, 0x00, 0x00)
MAGENTA = RGBColor(0xA9, 0x08, 0x8E)
HOT_PINK = RGBColor(0xFB, 0x47, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xE9, 0xE2, 0xE6)
MUTED = RGBColor(0xB6, 0xA8, 0xB0)
CARD = RGBColor(0x2C, 0x22, 0x27)
HAIRLINE = RGBColor(0x55, 0x33, 0x4A)

FONT_TITLE = 'Poppins'
FONT_BODY = 'Poppins'


def add_gradient_background(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    sp = bg.fill._xPr  # type: ignore[attr-defined]
    fill = sp.find(qn('a:solidFill'))
    if fill is not None:
        sp.remove(fill)
    grad_xml = (
        '<a:gradFill rotWithShape="1" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="000000"/></a:gs>'
        '<a:gs pos="55000"><a:srgbClr val="000000"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="34011F"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    sp.insert(2, etree.fromstring(grad_xml))
    bg.shadow.inherit = False
    sp_tree = bg._element.getparent()
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def add_glow_accents(slide) -> None:
    """Subtle accents to improve depth without changing theme."""
    orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.1), Inches(-0.7), Inches(3.2), Inches(3.2))
    orb1.fill.solid()
    orb1.fill.fore_color.rgb = MAGENTA
    orb1.fill.transparency = 0.82
    orb1.line.fill.background()
    orb1.shadow.inherit = False

    orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.8), Inches(6.1), Inches(2.2), Inches(2.2))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = HOT_PINK
    orb2.fill.transparency = 0.90
    orb2.line.fill.background()
    orb2.shadow.inherit = False


def add_face_caricature(slide, left, top, size=Inches(0.95), *, variant=0):
    """Add a small stylized face badge to make slides more vibrant."""
    skin_palette = [
        RGBColor(0xF0, 0xC8, 0xA0),
        RGBColor(0xD9, 0xA1, 0x7B),
        RGBColor(0xBF, 0x83, 0x5E),
    ]
    skin = skin_palette[variant % len(skin_palette)]

    # glow ring
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, left - Inches(0.05), top - Inches(0.05), size + Inches(0.1), size + Inches(0.1))
    ring.fill.solid()
    ring.fill.fore_color.rgb = HOT_PINK
    ring.fill.transparency = 0.65
    ring.line.fill.background()
    ring.shadow.inherit = False

    # face base
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    head.fill.solid()
    head.fill.fore_color.rgb = skin
    head.line.color.rgb = HAIRLINE
    head.line.width = Pt(0.8)
    head.shadow.inherit = False

    # hair cap
    hair = slide.shapes.add_shape(MSO_SHAPE.CLOUD, left + Inches(0.02), top - Inches(0.08), size - Inches(0.04), Inches(0.36))
    hair.fill.solid()
    hair.fill.fore_color.rgb = RGBColor(0x2A, 0x17, 0x23)
    hair.line.fill.background()
    hair.shadow.inherit = False

    # eyes
    eye_w = Inches(0.08)
    eye_h = Inches(0.08)
    le = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.26), top + Inches(0.40), eye_w, eye_h)
    re = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.61), top + Inches(0.40), eye_w, eye_h)
    for eye in (le, re):
        eye.fill.solid()
        eye.fill.fore_color.rgb = BLACK
        eye.line.fill.background()
        eye.shadow.inherit = False

    # mouth (use smiley overlay for easy editable expression)
    mouth = slide.shapes.add_shape(MSO_SHAPE.SMILEY_FACE, left + Inches(0.33), top + Inches(0.58), Inches(0.32), Inches(0.2))
    mouth.fill.solid()
    mouth.fill.fore_color.rgb = RGBColor(0xE0, 0x8A, 0xA0)
    mouth.fill.transparency = 0.3
    mouth.line.color.rgb = RGBColor(0x8C, 0x2F, 0x55)
    mouth.line.width = Pt(0.7)
    mouth.shadow.inherit = False


def add_text(slide, left, top, width, height, text: str, *, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font=FONT_BODY, italic=False, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.adjustments[0] = 0.06
    c.fill.solid()
    c.fill.fore_color.rgb = CARD
    c.line.color.rgb = HAIRLINE
    c.line.width = Pt(0.75)
    c.shadow.inherit = False
    return c


def add_pill(slide, left, top, label: str, fill=MAGENTA):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.6), Inches(0.42))
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.fill.background()
    add_text(slide, left, top + Inches(0.08), Inches(2.6), Inches(0.24), label, size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)


def add_footer(slide, text: str):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.95), Inches(12.1), Pt(1.0))
    line.fill.solid()
    line.fill.fore_color.rgb = HAIRLINE
    line.line.fill.background()
    add_text(slide, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.35), text, size=8.2, color=MUTED)


def set_notes(slide, notes: str):
    n = slide.notes_slide.notes_text_frame
    n.text = notes.strip()


def section_chrome(slide, eyebrow: str, title: str, no: str):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(9), Inches(0.32), eyebrow.upper(), size=11, bold=True,
             color=HOT_PINK, font=FONT_TITLE)
    add_text(slide, Inches(0.6), Inches(0.78), Inches(11.6), Inches(0.85), title, size=33, bold=True,
             color=WHITE, font=FONT_TITLE)
    hl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.62), Inches(2.2), Pt(2))
    hl.fill.solid(); hl.fill.fore_color.rgb = HOT_PINK; hl.line.fill.background()
    add_text(slide, Inches(11.6), Inches(0.45), Inches(1.2), Inches(0.3), no, size=11, bold=True,
             color=MUTED, align=PP_ALIGN.RIGHT, font=FONT_TITLE)


def icon_bullets(slide, left, top, width, items: Iterable[str], *, size=16):
    symbols = ('✦', '◆', '◉', '➤', '✧')
    y = top
    for idx, item in enumerate(items):
        sym = symbols[idx % len(symbols)]
        add_text(slide, left, y, width, Inches(0.55), f'{sym}  {item}', size=size, color=SOFT_WHITE)
        y += Inches(0.62)


def build_notes(title: str, bullets: Sequence[str], coach_line: str) -> str:
    detail_lines = '\n'.join(f'- {b}' for b in bullets)
    return (
        f'{coach_line}\n\n'
        f'Core beats for "{title}":\n'
        f'{detail_lines}\n\n'
        'Delivery cue: land one bullet per beat, pause briefly, then bridge forward.'
    )


def _chip_label(text: str) -> str:
    head = text.split(':', 1)[0].strip()
    if len(head) <= 26:
        return head
    words = text.split()
    return ' '.join(words[:4]).strip()


def add_bottom_chips(slide, bullets: Sequence[str]) -> None:
    """Use lower area for compact visual callouts."""
    chips = [_chip_label(b) for b in bullets[:3]]
    while len(chips) < 3:
        chips.append('Key point')

    y = Inches(5.45)
    for i, label in enumerate(chips):
        x = Inches(0.85 + i * 4.0)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(1.08))
        card.adjustments[0] = 0.16
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x24, 0x1B, 0x20)
        card.line.color.rgb = HAIRLINE
        card.line.width = Pt(0.7)
        card.shadow.inherit = False
        add_text(slide, x + Inches(0.18), y + Inches(0.17), Inches(3.25), Inches(0.28),
                 'KEY SIGNAL', size=9.5, bold=True, color=HOT_PINK, font=FONT_TITLE)
        add_text(slide, x + Inches(0.18), y + Inches(0.5), Inches(3.25), Inches(0.45),
                 label, size=13, bold=True, color=WHITE, font=FONT_TITLE)


def add_mini_bar_viz(slide, left, top, width, height, labels: Sequence[str], values: Sequence[float], *,
                     accent=HOT_PINK, title='DATA SNAPSHOT'):
    """Compact, reusable bar-chart card."""
    card = add_card(slide, left, top, width, height)
    _ = card
    add_text(slide, left + Inches(0.16), top + Inches(0.12), width - Inches(0.32), Inches(0.22),
             title, size=9.5, bold=True, color=HOT_PINK, font=FONT_TITLE)

    max_v = max(values) if values else 1
    n = max(1, len(labels))
    bar_area_w = width - Inches(0.5)
    slot = bar_area_w / n
    base_y = top + height - Inches(0.26)

    for i, (lab, val) in enumerate(zip(labels, values)):
        bar_h = Inches(0.9) * (val / max_v if max_v > 0 else 0.0)
        bx = left + Inches(0.25) + slot * i + Inches(0.08)
        bw = slot - Inches(0.22)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, base_y - bar_h, bw, bar_h)
        bar.adjustments[0] = 0.2
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent if i % 2 == 0 else MAGENTA
        bar.line.fill.background()
        bar.shadow.inherit = False
        add_text(slide, bx, base_y + Inches(0.03), bw, Inches(0.18), lab, size=8.5, color=MUTED,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text(slide, bx, base_y - bar_h - Inches(0.16), bw, Inches(0.14), f'{val:g}', size=8.5,
                 color=SOFT_WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)


def add_flow_viz(slide, left, top, width, *, steps: Sequence[str]):
    """Compact horizontal concept flow diagram."""
    n = max(1, len(steps))
    slot = width / n
    for i, step in enumerate(steps):
        x = left + slot * i + Inches(0.04)
        w = slot - Inches(0.08)
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, Inches(0.52))
        node.adjustments[0] = 0.25
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(0x24, 0x1B, 0x20)
        node.line.color.rgb = HOT_PINK if i == 0 else HAIRLINE
        node.line.width = Pt(0.75)
        node.shadow.inherit = False
        add_text(slide, x + Inches(0.08), top + Inches(0.13), w - Inches(0.16), Inches(0.24), step,
                 size=11, bold=True, color=SOFT_WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
        if i < n - 1:
            ax = x + w + Inches(0.015)
            arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, ax, top + Inches(0.18), Inches(0.12), Inches(0.16))
            arr.fill.solid()
            arr.fill.fore_color.rgb = HOT_PINK
            arr.line.fill.background()
            arr.shadow.inherit = False


def add_auto_visual(slide, title: str, bullets: Sequence[str], slide_index: int):
    """Inject concept/data visuals based on slide topic."""
    t = title.lower()
    if 'prompt engineering' in t:
        add_flow_viz(slide, Inches(0.92), Inches(4.78), Inches(11.0), steps=('Role', 'Context', 'Task', 'Verify'))
    elif 'gap' in t or 'evidence is thin' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('Real K-12', 'No classroom', 'No prompts'),
                         values=(2, 28, 18), title='STUDY LANDSCAPE')
    elif 'research question' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('RQ1', 'RQ2', 'RQ3'), values=(1, 1, 1), title='3-PILLAR DESIGN')
    elif 'hypotheses' in t:
        add_flow_viz(slide, Inches(0.92), Inches(4.78), Inches(11.0), steps=('H1', 'H2', 'H3', 'Falsifiers'))
    elif 'design' in t:
        add_flow_viz(slide, Inches(0.92), Inches(4.78), Inches(11.0), steps=('Pre', 'Module', 'Task', 'Post'))
    elif 'intervention' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('S1', 'S2', 'S3'), values=(6, 5, 5), title='RUBRIC DIMS / SCENARIO')
    elif 'rubric' in t or 'numbers' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('Rel', 'Clr', 'Con', 'Bkg', 'Elab', 'NoAns'),
                         values=(98, 85, 93, 96, 90, 88), title='PASS-FAIL AGREEMENT')
    elif 'validation' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('Target', 'Observed'),
                         values=(0.92, 0.95), title='KAPPA CHECK')
    elif 'argument' in t:
        add_flow_viz(slide, Inches(0.92), Inches(4.78), Inches(11.0), steps=('Evidence', 'Mechanism', 'Conclusion', 'Policy'))
    elif 'implications' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('PAUSD', 'Schools', 'Research'),
                         values=(3, 3, 3), title='IMPACT SURFACES')
    elif 'limitations' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('Sample', 'Power', 'Retention'),
                         values=(30, 1, 0.65), title='BOUNDARY CHECK')
    elif 'reflection' in t:
        add_flow_viz(slide, Inches(0.92), Inches(4.78), Inches(11.0), steps=('Start', 'Learn', 'Build', 'Now'))
    elif 'closing' in t or 'teach it' in t:
        add_mini_bar_viz(slide, Inches(9.15), Inches(2.02), Inches(3.55), Inches(2.95),
                         labels=('Use', 'Need', 'Scale'),
                         values=(72, 100, 1), title='WHY THIS MATTERS')


CITE = {
    'hook': 'Common Sense Media, 2025 • Gerlich, Societies, 2025',
    'modes': 'Lehmann et al., 2025 • Kosmyna et al., 2024 • Eltahir & Babiker, 2024',
    'debt': 'Kosmyna et al., 2024',
    'lever': 'MIT Sloan, 2025 • Dennison et al., 2024 • Park & Choo, 2024 • Chang et al., 2025',
    'gap': 'Chen et al., 2024 • Dennison et al., 2024 • Gogan, 2024',
    'rubric': 'Xiao et al., 2025',
}


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(blank)
    add_gradient_background(s)
    add_glow_accents(s)
    return s


# 1 Title
s = new_slide()
orb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.2), Inches(2.5), Inches(5.4), Inches(5.4))
orb.fill.solid(); orb.fill.fore_color.rgb = MAGENTA; orb.fill.transparency = 0.45; orb.line.fill.background()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.5), Inches(0.06), Inches(2.5))
bar.fill.solid(); bar.fill.fore_color.rgb = HOT_PINK; bar.line.fill.background()
add_text(s, Inches(0.95), Inches(1.15), Inches(8), Inches(0.35), 'AP RESEARCH • ORAL DEFENSE • 2025-26',
         size=12, bold=True, color=HOT_PINK, font=FONT_TITLE)
add_text(s, Inches(0.95), Inches(2.5), Inches(11.6), Inches(1.5), 'Teaching the Machine\nto Teach Us',
         size=58, bold=True, color=WHITE, font=FONT_TITLE, line_spacing=1.0)
add_text(s, Inches(0.95), Inches(5.1), Inches(11.2), Inches(0.8),
         'Prompt-engineering instruction, AI reasoning, and literacy in PAUSD', size=21, color=SOFT_WHITE)
add_text(s, Inches(0.95), Inches(6.45), Inches(11.2), Inches(0.3),
         'Jerry Yan • Palo Alto Unified School District • apr-sooty.vercel.app', size=12, color=MUTED)
set_notes(s, 'Open strong. One question: can a short lesson change how students think with AI?')

# 2 Hook
s = new_slide()
section_chrome(s, '01 • Why this matters', 'We are inside an AI paradigm shift', '2 / 24')
add_text(s, Inches(0.6), Inches(2.05), Inches(8.0), Inches(0.5), 'LLMs are already shaping student cognition.',
         size=24, bold=True, font=FONT_TITLE)
icon_bullets(s, Inches(0.6), Inches(2.85), Inches(8.1), [
    '72% of U.S. teens already use AI companions',
    'Heaviest users show the lowest critical-thinking scores',
    'Cognitive offloading and thinking scores: r = -0.75',
    'Without intervention, this scales into workforce habits',
], size=16)
add_card(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(4.8))
add_pill(s, Inches(9.25), Inches(2.2), 'STAT OF THE DAY')
add_text(s, Inches(9.25), Inches(2.85), Inches(3.3), Inches(1.6), '72%', size=82, bold=True, color=HOT_PINK,
         font=FONT_TITLE, line_spacing=1.0)
add_text(s, Inches(9.25), Inches(4.65), Inches(3.3), Inches(1.7), 'of U.S. teens use AI\nas a companion.', size=14)
add_face_caricature(s, Inches(8.95), Inches(5.6), size=Inches(0.82), variant=1)
add_footer(s, CITE['hook'])
set_notes(s, 'Three in four teens already use AI. The intervention question is urgent, not theoretical.')

# 3 Modes
s = new_slide()
section_chrome(s, '01 • Mechanism', 'Two modes of AI use', '3 / 24')
add_card(s, Inches(0.6), Inches(2.0), Inches(5.95), Inches(4.6))
add_pill(s, Inches(0.85), Inches(2.2), 'SUBSTITUTIVE', fill=RGBColor(0x8B, 0x14, 0x27))
add_text(s, Inches(0.85), Inches(2.85), Inches(5.3), Inches(0.5), 'AI replaces student thinking', size=18, bold=True)
icon_bullets(s, Inches(0.85), Inches(3.45), Inches(5.2), [
    'No attempt first', 'Lower understanding', 'Weaker EEG coupling', 'Higher dependence'
], size=14)
add_face_caricature(s, Inches(4.95), Inches(5.68), size=Inches(0.74), variant=0)
add_card(s, Inches(6.75), Inches(2.0), Inches(5.95), Inches(4.6))
add_pill(s, Inches(7.0), Inches(2.2), 'COMPLEMENTARY', fill=MAGENTA)
add_text(s, Inches(7.0), Inches(2.85), Inches(5.3), Inches(0.5), 'AI scaffolds student thinking', size=18, bold=True)
icon_bullets(s, Inches(7.0), Inches(3.45), Inches(5.2), [
    'Explain-first prompts', 'Higher comprehension gains', 'Prompt quality mediates effect', 'Better transfer'
], size=14)
add_face_caricature(s, Inches(11.15), Inches(5.68), size=Inches(0.74), variant=2)
add_footer(s, CITE['modes'])
set_notes(s, 'Contrast substitution vs complementarity clearly. Prompting skill is the dial between them.')

# 4 Cognitive debt
s = new_slide()
section_chrome(s, '01 • Neuroscience', 'Cognitive debt is measurable', '4 / 24')
add_text(s, Inches(0.6), Inches(2.1), Inches(7.8), Inches(1.2), 'MIT EEG (n = 55): LLM-assisted writing showed the weakest neural coupling.',
         size=20, bold=True, font=FONT_TITLE)
icon_bullets(s, Inches(0.6), Inches(3.4), Inches(7.8), [
    'Signal drop appears after one session',
    'Recall quality falls immediately after assisted writing',
    'Pattern aligns with offloading literature',
], size=15)
viz_img_candidates = [
    '/Users/yanzihao/.cursor/projects/Users-yanzihao-Downloads-AP-Research-Page/assets/Screenshot_2026-04-27_at_09.46.58-a1a5cb6d-7e0b-49e5-856c-5fd996b9c0f9.png',
    '/Users/yanzihao/.cursor/projects/Users-yanzihao-Downloads-AP-Research-Page/assets/Screenshot_2026-04-26_at_23.18.05-136bd56b-9821-45c3-b958-b13c4f51f8af.png',
]
viz_img = next((p for p in viz_img_candidates if Path(p).exists()), None)
if viz_img:
    add_card(s, Inches(0.6), Inches(4.78), Inches(8.2), Inches(2.1))
    s.shapes.add_picture(viz_img, Inches(0.78), Inches(4.95), width=Inches(7.84), height=Inches(1.72))
    add_text(s, Inches(0.82), Inches(6.72), Inches(7.8), Inches(0.16),
             'Concept visualization from cited EEG-style comparison graphic', size=8.5, color=MUTED)
add_card(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(4.6))
add_pill(s, Inches(9.25), Inches(2.2), 'MIT EEG')
add_text(s, Inches(9.25), Inches(2.9), Inches(3.3), Inches(1.3), '83.3%', size=66, bold=True, color=HOT_PINK,
         font=FONT_TITLE, line_spacing=1.0)
add_text(s, Inches(9.25), Inches(4.4), Inches(3.3), Inches(1.8), 'less unaided recall\nafter AI-assisted\nwriting.', size=13)
add_face_caricature(s, Inches(9.0), Inches(5.7), size=Inches(0.82), variant=2)
add_footer(s, CITE['debt'])
set_notes(s, 'Pause at 83.3%. Frame cognitive debt as compounding cost.')

# 5 Full image quote slide
s = new_slide()
img_candidates = [
    '/Users/yanzihao/.cursor/projects/Users-yanzihao-Downloads-AP-Research-Page/assets/Screenshot_2026-04-26_at_23.18.05-136bd56b-9821-45c3-b958-b13c4f51f8af.png',
    '/Users/yanzihao/.cursor/projects/Users-yanzihao-Downloads-AP-Research-Page/assets/Screenshot_2026-04-26_at_23.34.45-bc086a9d-a164-45c4-a57a-9b03dfc2098f.png',
]
img_path = next((p for p in img_candidates if Path(p).exists()), None)
if img_path:
    s.shapes.add_picture(img_path, 0, 0, width=SLIDE_W, height=SLIDE_H)
    ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    ov.fill.solid(); ov.fill.fore_color.rgb = BLACK; ov.fill.transparency = 0.45
    ov.line.fill.background()
add_text(s, Inches(1.0), Inches(2.6), Inches(11.0), Inches(1.8),
         'Literacy is the dial\nbetween AI harm and AI help.', size=52, bold=True, font=FONT_TITLE)
set_notes(s, 'This is your visual breath slide. One sentence only; let it land.')

# utility for compact section slides
slide_no = 6

def compact_slide(eyebrow, title, bullets, notes, cite=''):
    global slide_no
    ss = new_slide()
    section_chrome(ss, eyebrow, title, f'{slide_no} / 24')
    add_face_caricature(ss, Inches(11.65), Inches(1.95), variant=slide_no)
    bullet_list = list(bullets)
    icon_bullets(ss, Inches(0.85), Inches(2.05), Inches(11.2), bullet_list, size=18 if len(bullet_list) < 5 else 16)
    add_auto_visual(ss, title, bullet_list, slide_no)
    add_bottom_chips(ss, bullet_list)
    if cite:
        add_footer(ss, cite)
    set_notes(ss, build_notes(title, bullet_list, notes))
    slide_no += 1

compact_slide('01 • The lever', 'Prompt engineering = programmable language', [
    'Role: define who the model should be',
    'Context: include constraints and audience',
    'Task: specify artifact and format',
    'Verify: check claims before acceptance',
], 'Walk across Role, Context, Task, Verify. Keep this practical.', CITE['lever'])

compact_slide('01 • Gap', 'Curricula exist, evidence is thin', [
    'Strong frameworks already exist (CRAFT, Prompty, IDEA/PARTS)',
    'Few K-12 studies run controlled classroom comparisons',
    'Prompt transparency is often missing',
    'Restrictive-policy districts are understudied',
], 'Locate your contribution in one sentence: scored, real-student, controlled data.', CITE['gap'])

compact_slide('02 • Research question', 'One question, three measurable outcomes', [
    'RQ1: Substitutive-use behavior index',
    'RQ2: Prompt-quality rubric pass rates',
    'RQ3: Hallucination detect-and-correct accuracy',
], 'Read the RQ slowly; then tap each pillar.', '')

compact_slide('02 • Hypotheses', 'Directional predictions with falsifiers', [
    'H1: Treatment lowers substitutive-use index',
    'H2: Treatment raises rubric pass rates (largest in Background/Elaboration)',
    'H3: Treatment improves hallucination detection and correction',
    'Each hypothesis has an explicit null/falsifier test',
], 'State falsifiers out loud. That signals rigorous design.', '')

compact_slide('03 • Design', 'Two-arm asynchronous quasi-experiment', [
    'Volunteer sample; deterministic assignment via participant sequence',
    'Control: matched-length digital literacy reading',
    'Treatment: PE module + scenario tasks',
    'Total session time <= 60 minutes',
], 'Emphasize isolation: only one middle block differs.', '')
compact_slide('03 • Intervention', 'Three scenarios map to three RQs', [
    'S1 Ethical use -> RQ1 behavior',
    'S2 Iteration -> RQ2 prompt quality',
    'S3 Verification -> RQ3 hallucination handling',
    'Common scoring pipeline across all scenarios',
], 'The mapping is the spine. Keep this clean and fast.', '')

compact_slide('03 • Rubric', 'How prompts become numbers', [
    'Six binary dimensions (Xiao et al.)',
    'GPT-4o judge, temperature 0, JSON-only output',
    'Scenario-aware system prompt, fixed grading instructions',
    'Dimension-level pass rates feed inferential tests',
], 'Explain calibration choice and consistency benefits.', CITE['rubric'])

compact_slide('03 • Instruments', 'Pre/post measurement battery', [
    'Knowledge checks (T/F + open responses)',
    'Hallucination subtest with planted errors',
    'Confidence battery for placebo/familiarity control',
    'CRAFT task logs live behavior with AI',
], 'Behavioral logs are key: not just self-report.', '')

compact_slide('03 • Platform', 'Built end-to-end for replication', [
    'Next.js frontend + Supabase backend',
    'Typed API route for auto-grading',
    'SQL views flatten event logs for analysis',
    'Open-source workflow for district replication',
], 'Replication is a research value proposition, not a side note.', '')

compact_slide('04 • Validation', 'Can the model judge accurately?', [
    '15% stratified human audit sample',
    'Inter-rater target: kappa >= 0.80',
    'Auto-vs-gold target: kappa >= 0.92',
    'Preliminary overall agreement meets deployment threshold',
], 'Acknowledge weakest rubric dimension and ongoing calibration.', CITE['rubric'])

# result placeholders slides with mini chart boxes
for rid, title, hint in [
    ('RQ1', 'Results placeholder: Substitutive use', 'Insert bar chart + Welch t + Cohen d'),
    ('RQ2', 'Results placeholder: Prompt quality', 'Insert radar/line chart + mixed ANOVA'),
    ('RQ3', 'Results placeholder: Hallucination handling', 'Insert confusion grid + proportion tests'),
]:
    s = new_slide()
    section_chrome(s, '04 • Results', f'{rid} — {title}', f'{slide_no} / 24')
    add_card(s, Inches(0.6), Inches(2.0), Inches(8.2), Inches(4.7))
    add_text(s, Inches(0.9), Inches(2.3), Inches(7.6), Inches(0.5), '[ FIGURE PLACEHOLDER ]', size=14, bold=True,
             color=HOT_PINK, font=FONT_TITLE)
    add_text(s, Inches(0.9), Inches(3.05), Inches(7.6), Inches(1.0), hint, size=15)
    add_card(s, Inches(9.0), Inches(2.0), Inches(3.7), Inches(4.7))
    add_pill(s, Inches(9.25), Inches(2.2), 'SPEAK TRACK')
    add_face_caricature(s, Inches(11.55), Inches(2.15), size=Inches(0.82), variant=slide_no)
    add_text(s, Inches(9.25), Inches(2.9), Inches(3.3), Inches(3.3),
             'Direction\nEffect size\nUncertainty\nImplication', size=16, bold=True, color=SOFT_WHITE, line_spacing=1.35)
    set_notes(s, 'On defense day: fill with final stats. Lead with effect size, not just p-values.')
    slide_no += 1

compact_slide('05 • Argument', 'Evidence -> conclusion -> consequence', [
    'Evidence: substitution hurts, complementarity helps',
    'Mechanism: prompt skill shifts use mode',
    'Conclusion: short PE instruction can move behavior',
    'Consequence: literacy-first policy outperforms blanket bans',
], 'Run this as a logic chain, not as isolated claims.', '')

compact_slide('05 • Implications', 'What changes if findings hold', [
    'PAUSD: pilot short PE modules in ELA/history',
    'Schools: low-cost open-source intervention template',
    'Research: reusable scoring protocol + comparable metrics',
], 'Close each implication with a concrete implementation action.', '')

compact_slide('05 • Limitations', 'Boundaries and next studies', [
    'Single district and volunteer sampling',
    'Limited power for arm-specific effects',
    'Optional retention follow-up risks attrition bias',
    'Next: n ~ 90 multi-district replication',
], 'Explicit limits increase credibility. Keep tone confident.', '')

compact_slide('06 • Reflection', 'How my thinking evolved', [
    'From ban-first intuition to literacy-first framing',
    'From descriptive concern to measurable intervention',
    'From prompt tips to auditable rubric logic',
    'Current stance: teach the dial, not the fear',
], 'This slide earns reflection points. Be personal, but concise.', '')

compact_slide('07 • Closing', 'We do not ban literacy. We teach it.', [
    '72% are already using AI',
    'Your study tests a practical, replicable intervention',
    'Open stack + clear metrics make scaling realistic',
], 'Pause, then invite questions. End on agency.', '')

# 24 Works cited
s = new_slide()
section_chrome(s, '08 • Works cited', 'MLA sources supporting claims', '24 / 24')
works = [
    'Chang et al. (2025), npj Digital Medicine',
    'Chen et al. (2024), K-12 PE systematic review',
    'Common Sense Media (2025), teen AI use survey',
    'Dennison et al. (2024), Prompty / AAAI',
    'Eltahir & Babiker (2024), e-learning outcomes',
    'Gerlich (2025), critical thinking & offloading',
    'Gogan (2024), Gen/ReGen log',
    'Kosmyna et al. (2024), MIT EEG cognitive debt',
    'Lehmann et al. (2025), AI in classroom harm conditions',
    'MIT Sloan (2025), prompt essentials',
    'Park & Choo (2024), educator PE strategies',
    'Shojaee et al. (2025), reasoning limits',
    'Singhal et al. (2023), clinical knowledge in LLMs',
    'Xiao et al. (2025), rubric-anchored LLM-as-judge',
]
col1 = '\n'.join(f'• {w}' for w in works[:7])
col2 = '\n'.join(f'• {w}' for w in works[7:])
add_text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.9), col1, size=10, line_spacing=1.3, color=SOFT_WHITE)
add_text(s, Inches(6.75), Inches(2.0), Inches(6.0), Inches(4.9), col2, size=10, line_spacing=1.3, color=SOFT_WHITE)
set_notes(s, 'Keep on screen for panel verification. Do not read aloud.')

out_path = os.path.abspath('AP_Research_Presentation.pptx')
prs.save(out_path)
print(f'Saved {out_path} ({len(prs.slides)} slides)')
